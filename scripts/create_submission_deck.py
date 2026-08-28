"""
Generate the 5-slide submission PowerPoint deck.
Run: python scripts/create_submission_deck.py
Output: docs/Clinical_Document_Intelligence_Hub_Deck.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUT = PROJECT_ROOT / "docs" / "Clinical_Document_Intelligence_Hub_Deck_v2.pptx"

# Palette
NAVY = RGBColor(0x0A, 0x3D, 0x62)
TEAL = RGBColor(0x00, 0x96, 0x88)
TEAL_LIGHT = RGBColor(0xB2, 0xDF, 0xDB)
ACCENT = RGBColor(0x26, 0xA6, 0x9A)
DARK = RGBColor(0x26, 0x32, 0x38)
MID = RGBColor(0x54, 0x6E, 0x7A)
GRAY = RGBColor(0x90, 0xA4, 0xAE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFB, 0xFC)
CARD = RGBColor(0xF0, 0xF4, 0xF8)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def accent_stripe(slide) -> None:
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL
    stripe.line.fill.background()


def header(slide, title: str, slide_num: int) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(0), Inches(9.88), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.22), Inches(8.5), Inches(0.65))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    num = slide.shapes.add_textbox(Inches(9.0), Inches(0.3), Inches(0.6), Inches(0.4))
    np = num.text_frame.paragraphs[0]
    np.text = str(slide_num)
    np.font.size = Pt(14)
    np.font.color.rgb = TEAL_LIGHT
    np.alignment = PP_ALIGN.RIGHT


def footer(slide, text: str = "Clinical Document Intelligence Hub") -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.05), Inches(9.0), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_LIGHT
    line.line.fill.background()

    fb = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(9.0), Inches(0.3))
    fp = fb.text_frame.paragraphs[0]
    fp.text = text
    fp.font.size = Pt(9)
    fp.font.color.rgb = GRAY


def card(slide, left, top, width, height, title: str, body: str, accent=TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xCF, 0xD8, 0xDC)
    shape.line.width = Pt(0.75)

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.35), Inches(height - 0.25))
    tf = tb.text_frame
    tf.word_wrap = True
    tp = tf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(13)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    bp = tf.add_paragraph()
    bp.text = body
    bp.font.size = Pt(11)
    bp.font.color.rgb = MID
    bp.space_before = Pt(6)


def bullet_block(slide, left, top, width, height, title: str, items: list[str], title_color=NAVY) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True

    tp = tf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(14)
    tp.font.bold = True
    tp.font.color.rgb = title_color
    tp.space_after = Pt(10)

    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}" if not item.startswith("  ") else item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
        p.level = 0


def slide_1_problem(prs: Presentation) -> None:
    slide = blank_slide(prs)
    accent_stripe(slide)
    header(slide, "Problem Understanding and Objective", 1)

    card(
        slide, 0.55, 1.25, 4.25, 2.6,
        "The Problem",
        "Clinical and administrative staff spend hours manually reading intake forms, "
        "discharge summaries, lab reports, and physician notes to extract key information.",
    )
    card(
        slide, 5.0, 1.25, 4.45, 2.6,
        "Why It Matters",
        "Manual review is slow, inconsistent, and error-prone. Teams need structured, "
        "decision-ready outputs — not more document reading.",
        accent=RGBColor(0x5C, 0x6B, 0xC0),
    )

    obj_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.1), Inches(8.9), Inches(2.65))
    obj_box.fill.solid()
    obj_box.fill.fore_color.rgb = NAVY
    obj_box.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(8.4), Inches(2.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Our Objective"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEAL_LIGHT

    objectives = [
        "Accept clinical documents — text, PDF, and image",
        "Extract structured entities with evidence and confidence scores",
        "Generate patient summaries, risk flags, and next-step recommendations",
        "Merge multiple documents for the same patient into one unified view",
    ]
    for obj in objectives:
        op = tf.add_paragraph()
        op.text = f"  →  {obj}"
        op.font.size = Pt(12)
        op.font.color.rgb = WHITE
        op.space_before = Pt(5)

    footer(slide, "Synthetic demo data only · Decision-support POC, not a diagnostic system")


def slide_2_architecture(prs: Presentation) -> None:
    slide = blank_slide(prs)
    accent_stripe(slide)
    header(slide, "Solution Architecture and Design Flow", 2)

    steps = [
        ("1", "Upload", "TXT · PDF · Image"),
        ("2", "Ingest", "PyMuPDF · Tesseract"),
        ("3", "Extract", "Groq LLM API"),
        ("4", "Enrich", "Confidence · Normalize"),
        ("5", "Assess", "Rules · Knowledge"),
        ("6", "Output", "Summary · Dashboard"),
    ]

    x_start = 0.45
    box_w = 1.42
    gap = 0.08
    y = 1.35

    for i, (num, label, sub) in enumerate(steps):
        x = x_start + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(1.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL if i % 2 == 0 else NAVY
        shape.line.fill.background()

        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.text = num
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = TEAL_LIGHT
        p1.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(8)
        p3.font.color.rgb = TEAL_LIGHT
        p3.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(Inches(x + box_w), Inches(y + 0.45), Inches(gap + 0.05), Inches(0.4))
            ap = arrow.text_frame.paragraphs[0]
            ap.text = "›"
            ap.font.size = Pt(20)
            ap.font.bold = True
            ap.font.color.rgb = TEAL
            ap.alignment = PP_ALIGN.CENTER

    card(
        slide, 0.55, 2.95, 4.25, 1.55,
        "Multi-Document Merge",
        "Upload notes, discharge summaries, and lab reports for one patient. "
        "System deduplicates findings, tags provenance, and blocks merge on identity conflicts.",
    )
    card(
        slide, 5.0, 2.95, 4.45, 1.55,
        "Evaluation Layer",
        "Automated scoring against ground_truth.json and the public Superinsight "
        "benchmark — all metrics from live API calls.",
        accent=RGBColor(0x5C, 0x6B, 0xC0),
    )

    stack_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.75), Inches(8.9), Inches(0.85))
    stack_box.fill.solid()
    stack_box.fill.fore_color.rgb = CARD
    stack_box.line.color.rgb = TEAL_LIGHT

    st = slide.shapes.add_textbox(Inches(0.75), Inches(4.95), Inches(8.5), Inches(0.5))
    sp = st.text_frame.paragraphs[0]
    sp.text = "Tech Stack:  Python  ·  Groq LLM  ·  Pydantic  ·  Streamlit  ·  PyMuPDF  ·  Tesseract OCR"
    sp.font.size = Pt(12)
    sp.font.bold = True
    sp.font.color.rgb = NAVY
    sp.alignment = PP_ALIGN.CENTER

    ui_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.85), Inches(5.0), Inches(0.75))
    ui_box.fill.solid()
    ui_box.fill.fore_color.rgb = TEAL
    ui_box.line.fill.background()
    ut = ui_box.text_frame.paragraphs[0]
    ut.text = "Streamlit Dashboard  →  Patient Summary  ·  Workflow Flags  ·  JSON Export"
    ut.font.size = Pt(12)
    ut.font.bold = True
    ut.font.color.rgb = WHITE
    ut.alignment = PP_ALIGN.CENTER

    footer(slide)


def slide_3_highlights(prs: Presentation) -> None:
    slide = blank_slide(prs)
    accent_stripe(slide)
    header(slide, "Implementation Highlights", 3)

    features = [
        ("7-Step Pipeline", "Ingest → LLM extract → confidence score → terminology normalize → knowledge retrieval → rule assessment → summary generation"),
        ("Schema Validation", "Pydantic models enforce structure; every field carries an evidence quote from the source document"),
        ("Workflow Flags", "Transparent clinical rules — e.g. HbA1c > 9%, BP > 140/90 — surfaced as actionable alerts"),
        ("Multi-Doc Merge", "Combine physician notes, discharge PDFs, and lab images; dedupe by canonical name with source provenance"),
        ("Benchmark Eval", "28/28 fields correct on 3-case suite: internal Jane Doe + 2 Superinsight excerpts (live Groq metrics)"),
        ("Audit Export", "One-click JSON download of full extraction, assessments, and summary for downstream systems"),
    ]

    positions = [
        (0.55, 1.25), (5.0, 1.25),
        (0.55, 3.05), (5.0, 3.05),
        (0.55, 4.85), (5.0, 4.85),
    ]
    accents = [TEAL, NAVY, TEAL, NAVY, TEAL, RGBColor(0x5C, 0x6B, 0xC0)]

    for (left, top), (title, body), accent in zip(positions, features, accents):
        card(slide, left, top, 4.25 if left < 3 else 4.45, 1.55, title, body, accent=accent)

    footer(slide)


def slide_4_challenges(prs: Presentation) -> None:
    slide = blank_slide(prs)
    accent_stripe(slide)
    header(slide, "Challenges and Learnings", 4)

    rows, cols = 6, 2
    table = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.25), Inches(8.9), Inches(4.35)).table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(5.7)

    headers = ["Challenge", "Learning / Mitigation"]
    data = [
        ("OCR noise on scanned lab images", "Documented limitation; plain-text physician note is the most reliable demo path"),
        ("Groq token limits on long documents", "Evaluation uses curated excerpts; full documents available for UI walkthrough"),
        ("Intermittent LLM JSON validation errors", "Retry logic with exponential backoff in evaluation script"),
        ("Medical terminology variance across docs", "terminology_map.json normalizes aliases to canonical clinical names"),
        ("POC scope vs clinical deployment", "Clear UI disclaimer — decision support only, human review required"),
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

    for r, (chal, learn) in enumerate(data, start=1):
        for c, text in enumerate([chal, learn]):
            cell = table.cell(r, c)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.85), Inches(8.9), Inches(0.9))
    insight.fill.solid()
    insight.fill.fore_color.rgb = TEAL
    insight.line.fill.background()

    it = slide.shapes.add_textbox(Inches(0.75), Inches(6.05), Inches(8.5), Inches(0.55))
    ip = it.text_frame.paragraphs[0]
    ip.text = (
        "Key Takeaway:  Hybrid AI works best — LLM for flexible extraction, "
        "deterministic rules for clinical flags, human review for final decisions."
    )
    ip.font.size = Pt(12)
    ip.font.bold = True
    ip.font.color.rgb = WHITE
    ip.alignment = PP_ALIGN.CENTER

    footer(slide)


def slide_5_demo(prs: Presentation) -> None:
    slide = blank_slide(prs)
    accent_stripe(slide)
    header(slide, "Demo Summary and Next Steps", 5)

    card(
        slide, 0.55, 1.25, 4.25, 2.8,
        "What We Delivered",
        "• Streamlit prototype — single & multi-document modes\n"
        "• Synthetic Jane Doe diabetes demo case\n"
        "• Superinsight benchmark validation (Apache 2.0)\n"
        "• GitHub repo, README, and evaluation scripts",
    )
    card(
        slide, 5.0, 1.25, 4.45, 2.8,
        "Future Enhancements",
        "• FHIR / HL7 export for EHR integration\n"
        "• Clinician review queue for low-confidence fields\n"
        "• Production OCR (Azure Doc Intelligence)\n"
        "• Role-based access and audit logging",
        accent=RGBColor(0x5C, 0x6B, 0xC0),
    )

    link_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.3), Inches(8.9), Inches(1.5))
    link_box.fill.solid()
    link_box.fill.fore_color.rgb = NAVY
    link_box.line.fill.background()

    lt = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(8.4), Inches(1.2))
    tf = lt.text_frame
    tf.word_wrap = True

    lp = tf.paragraphs[0]
    lp.text = "Resources"
    lp.font.size = Pt(14)
    lp.font.bold = True
    lp.font.color.rgb = TEAL_LIGHT

    links = [
        "GitHub:  github.com/ayush2888/Clinical-Document-Intelligence",
        "Live Demo:  [Add Streamlit Cloud URL after deployment]",
        "Video Walkthrough:  [Attach 3-min screen recording if applicable]",
    ]
    for link in links:
        lp2 = tf.add_paragraph()
        lp2.text = f"  →  {link}"
        lp2.font.size = Pt(11)
        lp2.font.color.rgb = WHITE
        lp2.space_before = Pt(4)

    thanks = slide.shapes.add_textbox(Inches(0.55), Inches(6.2), Inches(8.9), Inches(0.6))
    tp = thanks.text_frame.paragraphs[0]
    tp.text = "Thank you."
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = TEAL
    tp.alignment = PP_ALIGN.CENTER

    footer(slide)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_1_problem(prs)
    slide_2_architecture(prs)
    slide_3_highlights(prs)
    slide_4_challenges(prs)
    slide_5_demo(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
